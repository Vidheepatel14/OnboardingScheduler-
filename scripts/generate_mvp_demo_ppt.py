from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(19, 41, 75)
ORANGE = RGBColor(242, 106, 33)
SLATE = RGBColor(69, 85, 102)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 41, 59)
GREEN = RGBColor(31, 122, 72)


BASE_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = BASE_DIR / "presentations"
OUTPUT_PATH = OUTPUT_DIR / "AI_Onboarding_Concierge_MVP_Demo.pptx"


def add_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, label: str, color: RGBColor = ORANGE) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.38),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.06), Inches(3.5), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), Inches(12.2), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(11)
    run.font.color.rgb = SLATE


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(0.75),
        Inches(1.7),
        Inches(0.4),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    accent.text_frame.text = "MVP DEMO"
    p = accent.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(8.8), Inches(1.4))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "AI Onboarding Concierge"
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(Inches(0.72), Inches(2.7), Inches(9.0), Inches(1.0))
    p = subtitle_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Streamlit + Hugging Face APIs + Tool Calling + RAG"
    run.font.name = "Aptos"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(226, 232, 240)

    p = subtitle_box.text_frame.add_paragraph()
    run = p.add_run()
    run.text = "DSBA 6010 | March 17, 2026"
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(203, 213, 225)

    callout = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.85),
        Inches(1.4),
        Inches(3.7),
        Inches(3.65),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(255, 255, 255)
    callout.line.fill.background()
    tf = callout.text_frame
    tf.word_wrap = True
    lines = [
        ("What the MVP proves", ORANGE, True),
        ("Grounded handbook answers", DARK, False),
        ("Task tracking and completion", DARK, False),
        ("Calendar-based onboarding scheduling", DARK, False),
        ("Uploaded PDF/image Q&A", DARK, False),
        ("Automatic HR escalation when docs fail", DARK, False),
    ]
    for idx, (text, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(9 if idx == 0 else 6)
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(18 if idx == 0 else 16)
        run.font.bold = bold
        run.font.color.rgb = color


def add_bullet_slide(prs: Presentation, section: str, title: str, bullets: list[str], highlight: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_top_bar(slide, section)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.7), Inches(7.7), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(7.6), Inches(4.9))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(14)
        p.bullet = True
        run = p.add_run()
        run.text = bullet
        run.font.name = "Aptos"
        run.font.size = Pt(19)
        run.font.color.rgb = DARK

    if highlight:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.8),
            Inches(1.55),
            Inches(3.7),
            Inches(3.0),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Presenter note"
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = ORANGE

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = highlight
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.color.rgb = DARK

    add_footer(slide, "AI Onboarding Concierge MVP")


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_top_bar(slide, "Architecture")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.7), Inches(9.0), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "How the MVP works"
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    boxes = [
        ("1. Streamlit UI", "Employee chats, sees tasks, uploads files"),
        ("2. AgentSession", "Routes the request and calls tools"),
        ("3. Tool Router", "Search policy, schedule, email, document QA"),
        ("4. Data + Services", "SQLite, Chroma, Google Calendar, Gmail"),
        ("5. Grounded Response", "Answer with sources or escalate to HR"),
    ]

    left = 0.55
    top = 1.75
    width = 2.35
    height = 1.65
    gap = 0.2

    for index, (heading, body) in enumerate(boxes):
        x = Inches(left + index * (width + gap))
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            Inches(top),
            Inches(width),
            Inches(height),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(226, 232, 240)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = heading
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = NAVY

        p = tf.add_paragraph()
        run = p.add_run()
        run.text = body
        run.font.name = "Aptos"
        run.font.size = Pt(14)
        run.font.color.rgb = DARK

    note = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(4.45),
        Inches(11.9),
        Inches(1.55),
    )
    note.fill.solid()
    note.fill.fore_color.rgb = RGBColor(255, 248, 240)
    note.line.color.rgb = RGBColor(253, 186, 116)
    tf = note.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Tech stack"
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = ORANGE
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = (
        "Hosted Hugging Face models for agent + policy answers, hybrid BM25/vector retrieval with reranking, "
        "SQLite for task state, and Google integrations for real scheduling/email actions."
    )
    run.font.name = "Aptos"
    run.font.size = Pt(16)
    run.font.color.rgb = DARK

    add_footer(slide, "Grounded answers + action-oriented tooling")


def add_demo_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_top_bar(slide, "Demo")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.7), Inches(9.0), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Recommended live demo flow"
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    steps = [
        ("1", "Tasks", "Show pending tasks and mark one complete."),
        ("2", "Policy Q&A", "Ask a PTO question and show grounded sources."),
        ("3", "Exception case", "Ask the photo / first-day exception question."),
        ("4", "Scheduling", "Check availability and schedule task IDs."),
        ("5", "Document QA", "Upload a PDF and ask a question about it."),
    ]

    for idx, (num, heading, body) in enumerate(steps):
        x = Inches(0.7 + (idx % 3) * 4.15)
        y = Inches(1.8 + (idx // 3) * 2.05)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            y,
            Inches(3.55),
            Inches(1.55),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)
        tf = card.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{num}. {heading}"
        run.font.name = "Aptos"
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.color.rgb = NAVY
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = body
        run.font.name = "Aptos"
        run.font.size = Pt(14)
        run.font.color.rgb = DARK

    add_footer(slide, "Warm up one policy query before presenting to cache the reranker")


def add_prompt_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_top_bar(slide, "Prompts")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.7), Inches(10.0), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Exact prompts to use in class"
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    prompts = [
        "What tasks are still pending for me?",
        "Are occasional and seasonal employees eligible for PTO?",
        "I am joining PNC's Customer Care Center. Do I need to upload my photo prior to my first day?",
        "When am I free for onboarding this week?",
        "Schedule task ids 16, 17, and 18 tomorrow within my available work hours.",
        "What does this uploaded PDF say about reimbursement?",
    ]

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(11.9), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, prompt in enumerate(prompts, start=1):
        p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = f"{idx}. {prompt}"
        run.font.name = "Courier New"
        run.font.size = Pt(18)
        run.font.color.rgb = DARK

    add_footer(slide, "Use broad scheduling windows to reduce demo risk")


def add_value_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT)
    add_top_bar(slide, "Value")

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.7), Inches(9.5), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Why this stands out"
    run.font.name = "Aptos Display"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.7),
        Inches(5.8),
        Inches(4.6),
    )
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = RGBColor(226, 232, 240)
    tf = left.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Employee benefit"
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = GREEN
    for text in [
        "One interface for questions, tasks, and scheduling",
        "Grounded answers with document sources instead of unsupported guesses",
        "Faster onboarding with fewer HR delays and fewer manual searches",
        "Can answer questions about uploaded PDFs and images, not just fixed FAQs",
    ]:
        p = tf.add_paragraph()
        p.bullet = True
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.color.rgb = DARK

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.95),
        Inches(1.7),
        Inches(5.6),
        Inches(4.6),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = RGBColor(226, 232, 240)
    tf = right.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Why it is more than a chatbot"
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = ORANGE
    for text in [
        "It takes actions through tools instead of only generating text",
        "It combines retrieval, task workflows, document QA, and escalation",
        "It uses hosted Hugging Face APIs, so no local LLM server is required",
        "It has a clear failure path: if the answer is not in documentation, escalate to HR",
    ]:
        p = tf.add_paragraph()
        p.bullet = True
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.color.rgb = DARK

    add_footer(slide, "Actionable, grounded, and demoable")


def add_close_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(8.8), Inches(1.1))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "MVP takeaways"
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.95), Inches(2.5), Inches(10.1), Inches(3.0))
    tf = body.text_frame
    for idx, text in enumerate(
        [
            "The app already demonstrates grounded handbook Q&A, task tracking, scheduling, document QA, and HR escalation.",
            "The strongest class demo is a short end-to-end workflow, not a deep technical explanation.",
            "Next work after the MVP: broader tests, better retrieval evaluation, and tighter production auth flows.",
        ]
    ):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = text
        run.font.name = "Aptos"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(226, 232, 240)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(6.15),
        Inches(2.4),
        Inches(0.45),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ORANGE
    tag.line.fill.background()
    tag.text_frame.text = "Questions"
    p = tag.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = WHITE


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "Problem",
        "Why this project exists",
        [
            "New hires often switch between PDFs, email, task trackers, and calendars just to finish onboarding.",
            "HR and managers spend time answering repetitive handbook questions that should be self-service.",
            "Scheduling onboarding tasks creates extra back-and-forth even when the information already exists.",
            "The MVP goal is a single assistant that can answer, act, and escalate when documentation is insufficient.",
        ],
        highlight="Frame this as a workflow problem, not only a chatbot problem.",
    )
    add_bullet_slide(
        prs,
        "Scope",
        "What the MVP can do today",
        [
            "Answer handbook and policy questions with retrieved document sources.",
            "Show pending onboarding tasks and mark tasks complete.",
            "Check calendar availability and schedule onboarding tasks.",
            "Answer questions about uploaded PDFs and images.",
            "Draft an HR escalation email when the documentation does not answer the question.",
        ],
        highlight="The strongest message is that the app combines knowledge + actions.",
    )
    add_architecture_slide(prs)
    add_demo_flow_slide(prs)
    add_prompt_slide(prs)
    add_value_slide(prs)
    add_close_slide(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(path)
